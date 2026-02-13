from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

slides_content = [
    "Welcome to My AI Video Demo",
    "This system converts PowerPoint to video",
    "Built with FastAPI, Docker, and FFmpeg"
]

for text in slides_content:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = text
    content.text = "This is sample content for testing."

prs.save("test.pptx")

print("Created test.pptx")
