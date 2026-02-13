from pptx import Presentation

def extract_slides(ppt_path):
    prs = Presentation(ppt_path)
    slides_data = []

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)

        slides_data.append({
            "slide_number": i,
            "text": " ".join(texts)
        })

    return slides_data
